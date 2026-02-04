"""
Unified Extraction Module for S-Core
Extracts text and images from all supported file formats:
- PDF, DOC/DOCX, PPT/PPTX, TXT
- Scanned documents, handwritten notes
"""

import io
import os
from typing import List, Tuple
from dataclasses import dataclass

import pdfplumber
from PIL import Image

# Optional imports with fallbacks
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ─────────────────────────────────────────────────────────────
# Data Classes
# ─────────────────────────────────────────────────────────────

@dataclass
class ExtractedContent:
    """Container for extracted content from a file"""
    text_chunks: List[str]
    images: List[bytes]  # Raw image bytes
    source_file: str

# ─────────────────────────────────────────────────────────────
# PDF Extraction
# ─────────────────────────────────────────────────────────────

def extract_pdf(path: str) -> ExtractedContent:
    """
    Extract text and images from PDF.
    Handles both text-based and scanned PDFs.
    """
    text_chunks = []
    images = []
    
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages):
            # Extract text
            text = page.extract_text()
            if text and text.strip():
                text_chunks.append(text.strip())
            
            # Extract images from page
            if hasattr(page, 'images') and page.images:
                for img_info in page.images:
                    try:
                        # Get image bbox and extract
                        x0 = img_info.get('x0', 0)
                        y0 = img_info.get('top', 0)
                        x1 = img_info.get('x1', page.width)
                        y1 = img_info.get('bottom', page.height)
                        
                        # Crop page to image area and convert to PIL Image
                        cropped = page.within_bbox((x0, y0, x1, y1))
                        img = cropped.to_image(resolution=150)
                        
                        # Convert to bytes
                        img_buffer = io.BytesIO()
                        img.save(img_buffer, format='PNG')
                        images.append(img_buffer.getvalue())
                    except Exception:
                        # Skip problematic images
                        continue
    
    return ExtractedContent(
        text_chunks=text_chunks,
        images=images,
        source_file=os.path.basename(path)
    )

# ─────────────────────────────────────────────────────────────
# Word Document Extraction
# ─────────────────────────────────────────────────────────────

def extract_docx(path: str) -> ExtractedContent:
    """Extract text and images from Word documents"""
    if not HAS_DOCX:
        raise ImportError("python-docx not installed. Run: pip install python-docx")
    
    text_chunks = []
    images = []
    
    doc = DocxDocument(path)
    
    # Extract text from paragraphs
    current_chunk = []
    for para in doc.paragraphs:
        if para.text.strip():
            current_chunk.append(para.text.strip())
        elif current_chunk:
            text_chunks.append('\n'.join(current_chunk))
            current_chunk = []
    
    if current_chunk:
        text_chunks.append('\n'.join(current_chunk))
    
    # Extract text from tables
    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                table_text.append(' | '.join(row_text))
        if table_text:
            text_chunks.append('\n'.join(table_text))
    
    # Extract images
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                images.append(image_data)
            except Exception:
                continue
    
    return ExtractedContent(
        text_chunks=text_chunks,
        images=images,
        source_file=os.path.basename(path)
    )

# ─────────────────────────────────────────────────────────────
# PowerPoint Extraction
# ─────────────────────────────────────────────────────────────

def extract_pptx(path: str) -> ExtractedContent:
    """Extract text and images from PowerPoint presentations"""
    if not HAS_PPTX:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    
    text_chunks = []
    images = []
    
    prs = Presentation(path)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text.append(' | '.join(row_text))
            
            # Extract images
            if hasattr(shape, "image"):
                try:
                    images.append(shape.image.blob)
                except Exception:
                    continue
        
        if slide_text:
            text_chunks.append(f"[Slide {slide_num}]\n" + '\n'.join(slide_text))
    
    return ExtractedContent(
        text_chunks=text_chunks,
        images=images,
        source_file=os.path.basename(path)
    )

# ─────────────────────────────────────────────────────────────
# Plain Text Extraction
# ─────────────────────────────────────────────────────────────

def extract_txt(path: str) -> ExtractedContent:
    """Extract text from plain text files"""
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    # Split into paragraphs
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    return ExtractedContent(
        text_chunks=paragraphs if paragraphs else [content],
        images=[],
        source_file=os.path.basename(path)
    )

# ─────────────────────────────────────────────────────────────
# Image File Extraction
# ─────────────────────────────────────────────────────────────

def extract_image(path: str) -> ExtractedContent:
    """Load image file (for scanned notes, handwritten pages)"""
    with open(path, 'rb') as f:
        image_data = f.read()
    
    return ExtractedContent(
        text_chunks=[],
        images=[image_data],
        source_file=os.path.basename(path)
    )

# ─────────────────────────────────────────────────────────────
# Unified Extraction Interface
# ─────────────────────────────────────────────────────────────

EXTRACTORS = {
    'pdf': extract_pdf,
    'docx': extract_docx,
    'doc': extract_docx,  # Will need conversion in practice
    'pptx': extract_pptx,
    'ppt': extract_pptx,   # Will need conversion in practice
    'txt': extract_txt,
    'md': extract_txt,
    'png': extract_image,
    'jpg': extract_image,
    'jpeg': extract_image,
    'gif': extract_image,
    'bmp': extract_image,
    'tiff': extract_image,
    'tif': extract_image,
}

def extract_file(path: str) -> ExtractedContent:
    """
    Universal file extraction.
    Automatically detects file type and extracts content.
    """
    ext = path.lower().split('.')[-1]
    
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file type: .{ext}")
    
    return EXTRACTORS[ext](path)

def get_supported_extensions() -> List[str]:
    """Get list of supported file extensions"""
    return list(EXTRACTORS.keys())
